from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # blank

# Add slide title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "RAG System Workflow"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)

# Colors
ORANGE_LINE = (255, 140, 0)
ORANGE_FILL = (255, 248, 220)
GREEN_LINE = (50, 205, 50)
GREEN_FILL = (240, 255, 240)
BLUE_LINE = (65, 105, 225)
BLUE_FILL = (240, 248, 255)
PURPLE_LINE = (138, 43, 226)
PURPLE_FILL = (248, 248, 255)

def add_box(slide, text, left, top, width, height, bg_color, line_color, line_dash=False, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)
    
    line = shape.line
    line.color.rgb = RGBColor(*line_color)
    line.width = Pt(2)
    if line_dash:
        line.dash_style = 4
        
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
    return shape

def add_simple_arrow(slide, x, y, width, height=0.12):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y - height/2), Inches(width), Inches(height))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(150, 150, 150)
    shape.line.fill.background()

def add_label(slide, text, left, top, font_size=12, bold=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(8), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(0, 0, 0)

# SECTION 1: BUILD
add_label(slide1, "① BUILD — Prepare dataset (python build_index.py)", 0.5, 0.7, 13)

add_box(slide1, "Knowledge Base\ndata/animal_facts_qa.txt\nQ&A pairs", 0.5, 1.2, 1.6, 1.0, BLUE_FILL, BLUE_LINE, font_size=10)
add_simple_arrow(slide1, 2.1, 1.7, 0.3)
add_box(slide1, "document_loader\nload_qa_file()", 2.4, 1.2, 1.6, 1.0, BLUE_FILL, BLUE_LINE, font_size=10)
add_simple_arrow(slide1, 4.0, 1.7, 0.3)
add_box(slide1, "text_splitter\nbuild_chunks()\nSize: 400, Overlap: 50", 4.3, 1.2, 1.6, 1.0, BLUE_FILL, BLUE_LINE, font_size=10)
add_simple_arrow(slide1, 5.9, 1.7, 0.3)
add_box(slide1, "embedding_model\nparaphrase-multilingual-MiniLM-L12-v2", 6.2, 1.2, 1.6, 1.0, BLUE_FILL, BLUE_LINE, font_size=10)

add_simple_arrow(slide1, 7.8, 1.4, 0.3)
add_box(slide1, "vector_store\nFAISS index", 8.1, 0.9, 1.6, 0.7, BLUE_FILL, BLUE_LINE, font_size=10)

add_simple_arrow(slide1, 7.8, 2.0, 0.3)
add_box(slide1, "build_bm25\nBuilt in hybrid_retriever", 8.1, 1.7, 1.6, 0.7, BLUE_FILL, BLUE_LINE, font_size=10)

add_simple_arrow(slide1, 9.7, 1.4, 0.3)
add_box(slide1, "vector_db/\ndocument.index\nbm25_index.pkl\nchunk_store.json", 10.0, 0.9, 2.0, 1.5, PURPLE_FILL, PURPLE_LINE, font_size=10)

add_box(slide1, "index_meta.py\nSave dataset metadata", 8.1, 2.5, 1.6, 0.4, PURPLE_FILL, PURPLE_LINE, font_size=10)
add_simple_arrow(slide1, 9.7, 2.7, 0.3)

# Add line separator 1
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
shape.line.fill.background()

# SECTION 2: QUERY
add_label(slide1, "② QUERY — Answer questions (python main.py)", 0.5, 3.1, 13)

add_box(slide1, "[INPUT]\nUser Question\ninput()", 0.5, 3.6, 1.6, 1.0, GREEN_FILL, GREEN_LINE, font_size=10)
add_simple_arrow(slide1, 2.1, 4.1, 0.25)
add_box(slide1, "query_transform\nRewrite • multi • HyDE\n(OFF by default)", 2.35, 3.6, 1.6, 1.0, GREEN_FILL, GREEN_LINE, line_dash=True, font_size=10)
add_simple_arrow(slide1, 3.95, 4.1, 0.25)
add_box(slide1, "[RETRIEVAL]\nhybrid_retriever\nBM25 + dense → RRF\nTop 20 chunks", 4.2, 3.6, 1.8, 1.0, ORANGE_FILL, ORANGE_LINE, font_size=10)
add_simple_arrow(slide1, 6.0, 4.1, 0.25)
add_box(slide1, "rerankers\ncross-encoder: 20→3\n(OFF by default)", 6.25, 3.6, 1.6, 1.0, GREEN_FILL, GREEN_LINE, line_dash=True, font_size=10)
add_simple_arrow(slide1, 7.85, 4.1, 0.25)

add_box(slide1, "[CONTEXT]\nprompt_templates\nformat_context()\nbuild_messages()", 8.1, 3.6, 1.6, 1.0, ORANGE_FILL, ORANGE_LINE, font_size=10)
add_simple_arrow(slide1, 9.7, 4.1, 0.25)

add_box(slide1, "[LLM]\ngenerator\nGenerate answer via LLM\n(ollama/openai/gemini)", 9.95, 3.6, 1.6, 1.0, ORANGE_FILL, ORANGE_LINE, font_size=10)
add_simple_arrow(slide1, 11.55, 4.1, 0.25)

add_box(slide1, "[OUTPUT]\nAnswer\nReturn response +\ncitations [1] [2]", 11.8, 3.6, 1.5, 1.0, GREEN_FILL, GREEN_LINE, font_size=10)

# Memory box
add_box(slide1, "memory - last 6 turns\nKeep conversation history", 7.0, 4.8, 2.5, 0.5, ORANGE_FILL, ORANGE_LINE, font_size=10)

# Main.py loads index
add_box(slide1, "main.py loads index_meta\nCheck dataset & build new index if needed", 0.5, 4.8, 2.5, 0.5, GREEN_FILL, GREEN_LINE, font_size=10)

# Add line separator 2
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
shape.line.fill.background()

# SECTION 3: EVALUATION
add_label(slide1, "③ EVALUATION — Measure and improve", 0.5, 5.5, 13)

add_box(slide1, "golden_set.json\n60 questions\n140 answers", 0.5, 6.0, 1.6, 1.1, BLUE_FILL, BLUE_LINE, font_size=10)
add_simple_arrow(slide1, 2.1, 6.275, 0.25)
add_simple_arrow(slide1, 2.1, 6.825, 0.25)

add_box(slide1, "eval_retrieval.py", 2.35, 6.05, 1.6, 0.45, BLUE_FILL, BLUE_LINE, font_size=10)
add_simple_arrow(slide1, 3.95, 6.275, 0.25)
add_box(slide1, "Retrieval Report\noutputs/eval_retrieval.json", 4.2, 6.05, 2.0, 0.45, PURPLE_FILL, PURPLE_LINE, font_size=9)
add_simple_arrow(slide1, 6.2, 6.275, 0.25)
add_box(slide1, "metrics.py\nHit@k / MRR / nDCG", 6.45, 6.05, 1.8, 0.45, PURPLE_FILL, PURPLE_LINE, font_size=10)

add_box(slide1, "eval_generation.py", 2.35, 6.6, 1.6, 0.45, BLUE_FILL, BLUE_LINE, font_size=10)
add_simple_arrow(slide1, 3.95, 6.825, 0.25)
add_box(slide1, "Generation Report\noutputs/eval_generation.json", 4.2, 6.6, 2.0, 0.45, PURPLE_FILL, PURPLE_LINE, font_size=9)
add_simple_arrow(slide1, 6.2, 6.825, 0.25)
add_box(slide1, "Compare Configs\n(dense/bm25/hybrid)", 6.45, 6.6, 1.8, 0.45, PURPLE_FILL, PURPLE_LINE, font_size=10)

# Legend (moved to right side)
add_label(slide1, "LEGEND (Configurable Flags in config.py):", 8.5, 5.5, 10)
add_box(slide1, "Always ON", 8.5, 6.0, 1.3, 0.4, ORANGE_FILL, ORANGE_LINE, font_size=9)
add_box(slide1, "Configurable (ON)", 10.0, 6.0, 1.4, 0.4, GREEN_FILL, GREEN_LINE, font_size=9)
add_box(slide1, "Disabled by default", 11.6, 6.0, 1.4, 0.4, GREEN_FILL, GREEN_LINE, line_dash=True, font_size=9)
add_box(slide1, "Data / Resources", 8.5, 6.6, 1.3, 0.4, BLUE_FILL, BLUE_LINE, font_size=9)
add_box(slide1, "Outputs / Reports", 10.0, 6.6, 1.4, 0.4, PURPLE_FILL, PURPLE_LINE, font_size=9)


# SLIDE 2
slide2 = prs.slides.add_slide(prs.slide_layouts[1]) # title and content
title, content = slide2.shapes.title, slide2.placeholders[1]
title.text = "RAG System Workflow Explanation"

tf = content.text_frame
tf.text = "1. Input"
p = tf.paragraphs[0]
p.font.bold = True
p.level = 0

p = tf.add_paragraph()
p.text = "User inputs a question via the terminal prompt in main.py."
p.level = 1
p = tf.add_paragraph()
p.text = "The query is normalized and transformed (if enabled) via query_transform.py."
p.level = 1

p = tf.add_paragraph()
p.text = "2. Retrieval"
p.font.bold = True
p.level = 0

p = tf.add_paragraph()
p.text = "Queries are sent to HybridRetriever to search the indexed chunks (FAISS + BM25)."
p.level = 1

p = tf.add_paragraph()
p.text = "3. Context"
p.font.bold = True
p.level = 0

p = tf.add_paragraph()
p.text = "The top retrieved chunks are formatted into numbered references by prompt_templates.py."
p.level = 1

p = tf.add_paragraph()
p.text = "4. LLM"
p.font.bold = True
p.level = 0

p = tf.add_paragraph()
p.text = "The constructed prompt is sent to the LLM (Ollama, OpenAI, or Gemini) via generator.py."
p.level = 1

p = tf.add_paragraph()
p.text = "5. Output"
p.font.bold = True
p.level = 0

p = tf.add_paragraph()
p.text = "The system appends a disclaimer, caches the conversation, and returns the response."
p.level = 1

p = tf.add_paragraph()
p.text = "6. Evaluation (New)"
p.font.bold = True
p.level = 0

p = tf.add_paragraph()
p.text = "Uses a golden_set.json to measure retrieval accuracy (Hit@k, MRR) and generation quality."
p.level = 1

# Reduce font size of bullets to fit well
for p in tf.paragraphs:
    if p.level == 1:
        p.font.size = Pt(14)
    else:
        p.font.size = Pt(16)

prs.save("RAG_Workflow_Presentation_v2.pptx")
